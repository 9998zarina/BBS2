#!/usr/bin/env python3
"""
BBS AI Assessment System - Sales Presentation Generator
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9
prs.slide_height = Inches(7.5)

# Color scheme
PURPLE = RGBColor(139, 92, 246)  # Purple
BLUE = RGBColor(59, 130, 246)    # Blue
DARK = RGBColor(31, 41, 55)      # Dark gray
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(243, 244, 246)
GREEN = RGBColor(34, 197, 94)
RED = RGBColor(239, 68, 68)

def add_title_slide(title, subtitle=""):
    """Add a title slide with gradient-like background"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background shape (purple)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PURPLE
    bg.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_section_slide(title):
    """Add a section divider slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BLUE
    bg.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(title, content_items, two_column=False):
    """Add a content slide with bullet points"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # White background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    # Purple header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PURPLE
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if two_column and len(content_items) > 1:
        # Left column
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(6), Inches(5.5))
        tf = left_box.text_frame
        tf.word_wrap = True

        for i, item in enumerate(content_items[0]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(20)
            p.font.color.rgb = DARK
            p.space_after = Pt(12)

        # Right column
        right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.6), Inches(6), Inches(5.5))
        tf = right_box.text_frame
        tf.word_wrap = True

        for i, item in enumerate(content_items[1]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(20)
            p.font.color.rgb = DARK
            p.space_after = Pt(12)
    else:
        # Single column
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12.333), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, item in enumerate(content_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(24)
            p.font.color.rgb = DARK
            p.space_after = Pt(16)

    return slide

def add_comparison_slide(title, left_title, left_items, right_title, right_items):
    """Add a comparison slide with two columns"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # White background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    # Purple header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PURPLE
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Left column header (red/gray)
    left_header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(5.8), Inches(0.7))
    left_header.fill.solid()
    left_header.fill.fore_color.rgb = RGBColor(220, 220, 220)
    left_header.line.fill.background()

    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.55), Inches(5.8), Inches(0.7))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK
    p.alignment = PP_ALIGN.CENTER

    # Right column header (green/blue)
    right_header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.8), Inches(0.7))
    right_header.fill.solid()
    right_header.fill.fore_color.rgb = PURPLE
    right_header.line.fill.background()

    right_title_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.55), Inches(5.8), Inches(0.7))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Left column content
    left_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.4), Inches(5.5), Inches(4.5))
    tf = left_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"✗ {item}"
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(150, 150, 150)
        p.space_after = Pt(14)

    # Right column content
    right_box = slide.shapes.add_textbox(Inches(7), Inches(2.4), Inches(5.5), Inches(4.5))
    tf = right_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"✓ {item}"
        p.font.size = Pt(20)
        p.font.color.rgb = DARK
        p.space_after = Pt(14)

    return slide

def add_feature_slide(title, features):
    """Add a feature showcase slide with icons"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # White background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    # Purple header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PURPLE
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Feature boxes (2x2 or 3x2 grid)
    cols = 3 if len(features) > 4 else 2
    box_width = Inches(3.8) if cols == 3 else Inches(5.5)

    for i, (feat_title, feat_desc) in enumerate(features):
        row = i // cols
        col = i % cols

        x = Inches(0.5) + col * (box_width + Inches(0.3))
        y = Inches(1.6) + row * Inches(2.8)

        # Feature box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_width, Inches(2.5))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.fill.background()

        # Feature title
        ft_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.3), box_width - Inches(0.4), Inches(0.5))
        tf = ft_box.text_frame
        p = tf.paragraphs[0]
        p.text = feat_title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = PURPLE

        # Feature description
        fd_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.9), box_width - Inches(0.4), Inches(1.4))
        tf = fd_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = feat_desc
        p.font.size = Pt(16)
        p.font.color.rgb = DARK

    return slide

def add_table_slide(title, headers, rows):
    """Add a slide with a table"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # White background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    # Purple header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PURPLE
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Table
    num_rows = len(rows) + 1
    num_cols = len(headers)

    table_width = Inches(12)
    table_height = Inches(0.5) * num_rows

    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.667), Inches(1.6), table_width, table_height).table

    # Set column widths
    col_width = table_width / num_cols
    for i in range(num_cols):
        table.columns[i].width = int(col_width)

    # Header row
    for i, header_text in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = PURPLE
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # Data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(16)
            p.font.color.rgb = DARK
            p.alignment = PP_ALIGN.CENTER

    return slide

def add_process_slide(title, steps):
    """Add a process flow slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # White background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    # Purple header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PURPLE
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Process steps
    num_steps = len(steps)
    step_width = Inches(2.2)
    total_width = step_width * num_steps + Inches(0.5) * (num_steps - 1)
    start_x = (prs.slide_width - total_width) / 2

    for i, (step_num, step_title, step_desc) in enumerate(steps):
        x = start_x + i * (step_width + Inches(0.5))
        y = Inches(2)

        # Circle with number
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.6), y, Inches(1), Inches(1))
        circle.fill.solid()
        circle.fill.fore_color.rgb = PURPLE
        circle.line.fill.background()

        num_box = slide.shapes.add_textbox(x + Inches(0.6), y + Inches(0.2), Inches(1), Inches(0.6))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(step_num)
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Arrow (except last)
        if i < num_steps - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + step_width, y + Inches(0.3), Inches(0.5), Inches(0.4))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(200, 200, 200)
            arrow.line.fill.background()

        # Step title
        st_box = slide.shapes.add_textbox(x, y + Inches(1.2), step_width, Inches(0.5))
        tf = st_box.text_frame
        p = tf.paragraphs[0]
        p.text = step_title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = DARK
        p.alignment = PP_ALIGN.CENTER

        # Step description
        sd_box = slide.shapes.add_textbox(x, y + Inches(1.7), step_width, Inches(2))
        tf = sd_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = step_desc
        p.font.size = Pt(14)
        p.font.color.rgb = DARK
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_stats_slide(title, stats):
    """Add a statistics highlight slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Light purple background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(245, 243, 255)
    bg.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PURPLE
    p.alignment = PP_ALIGN.CENTER

    # Stats boxes
    num_stats = len(stats)
    box_width = Inches(3.5)
    total_width = box_width * num_stats + Inches(0.5) * (num_stats - 1)
    start_x = (prs.slide_width - total_width) / 2

    for i, (stat_value, stat_label) in enumerate(stats):
        x = start_x + i * (box_width + Inches(0.5))
        y = Inches(2.5)

        # Stat box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_width, Inches(3.5))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = PURPLE
        box.line.width = Pt(2)

        # Stat value
        sv_box = slide.shapes.add_textbox(x, y + Inches(0.5), box_width, Inches(1.5))
        tf = sv_box.text_frame
        p = tf.paragraphs[0]
        p.text = stat_value
        p.font.size = Pt(56)
        p.font.bold = True
        p.font.color.rgb = PURPLE
        p.alignment = PP_ALIGN.CENTER

        # Stat label
        sl_box = slide.shapes.add_textbox(x, y + Inches(2), box_width, Inches(1))
        tf = sl_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = stat_label
        p.font.size = Pt(20)
        p.font.color.rgb = DARK
        p.alignment = PP_ALIGN.CENTER

    return slide

# ============================================
# CREATE PRESENTATION
# ============================================

# Slide 1: Title
add_title_slide(
    "BBS AI Assessment System",
    "AI 기반 Berg Balance Scale 자동 평가 시스템"
)

# Slide 2: Problem Statement
add_comparison_slide(
    "기존 BBS 평가의 문제점",
    "기존 방식",
    [
        "수동 관찰 및 기록 (30분+ 소요)",
        "평가자 간 편차 발생",
        "실시간 관찰 필요",
        "경험에 의존한 주관적 판단",
        "데이터 정량화 어려움"
    ],
    "AI 솔루션",
    [
        "자동 영상 분석 (5분 내 완료)",
        "일관된 객관적 평가",
        "영상 기반 사후 분석 가능",
        "데이터 기반 정량적 측정",
        "상세 분석 리포트 제공"
    ]
)

# Slide 3: Key Stats
add_stats_slide(
    "핵심 성과 지표",
    [
        ("5분", "평가 소요 시간\n(기존 30분)"),
        ("14개", "자동 분석\nBBS 검사 항목"),
        ("85%+", "예상 채점\n정확도")
    ]
)

# Slide 4: How it works
add_process_slide(
    "작동 원리",
    [
        (1, "영상 업로드", "정면/측면\n영상 업로드"),
        (2, "AI 포즈 분석", "YOLOv8로\n17개 관절 추적"),
        (3, "검사 자동 인식", "동작 패턴으로\n검사 종류 판별"),
        (4, "자동 채점", "BBS 기준\n0-4점 채점"),
        (5, "결과 리포트", "상세 분석\n결과 제공")
    ]
)

# Slide 5: Section - BBS 검사 소개
add_section_slide("BBS (Berg Balance Scale) 검사")

# Slide 6: BBS Overview
add_content_slide(
    "BBS 검사란?",
    [
        "노인 및 신경계 환자의 균형 능력을 평가하는 국제 표준 검사",
        "총 14개 검사 항목으로 구성",
        "각 항목 0-4점, 총점 56점 만점",
        "낙상 위험도 예측에 활용",
        "물리치료사, 재활의학과에서 광범위하게 사용"
    ]
)

# Slide 7: Score Interpretation
add_table_slide(
    "점수 해석 기준",
    ["점수 범위", "낙상 위험도", "보행 보조기구", "권장 조치"],
    [
        ["41-56점", "낮음 (Low)", "불필요", "유지 관리"],
        ["21-40점", "중간 (Medium)", "필요할 수 있음", "재활 치료 권장"],
        ["0-20점", "높음 (High)", "필요", "집중 재활 필요"]
    ]
)

# Slide 8: 14 Test Items (1-7)
add_table_slide(
    "14개 검사 항목 (1-7)",
    ["번호", "검사 항목", "주요 측정", "시간 기준"],
    [
        ["1", "앉아서 일어서기", "손 사용, 안정성", "-"],
        ["2", "지지 없이 서기", "CoM 안정성", "2분"],
        ["3", "지지 없이 앉기", "등 떼고 앉기", "2분"],
        ["4", "서서 앉기", "하강 제어", "-"],
        ["5", "이동하기", "안전한 이동", "-"],
        ["6", "눈 감고 서기", "균형 유지", "10초"],
        ["7", "두 발 모으고 서기", "좌우/앞뒤 흔들림", "1분"]
    ]
)

# Slide 9: 14 Test Items (8-14)
add_table_slide(
    "14개 검사 항목 (8-14)",
    ["번호", "검사 항목", "주요 측정", "시간 기준"],
    [
        ["8", "팔 뻗어 내밀기", "손목 이동 거리", "-"],
        ["9", "바닥 물건 집기", "몸통 굴곡", "-"],
        ["10", "뒤돌아보기", "척추 회전", "-"],
        ["11", "360도 회전", "스텝 수, 시간", "4초"],
        ["12", "발판 발 교대", "스텝 수", "20초 8회"],
        ["13", "일렬로 서기", "tandem 자세", "30초"],
        ["14", "한 발로 서기", "단일 다리 균형", "10초"]
    ]
)

# Slide 10: Section - 핵심 기능
add_section_slide("핵심 기능")

# Slide 11: Key Features
add_feature_slide(
    "주요 기능",
    [
        ("자동 검사 인식", "영상만 업로드하면 AI가 14개 검사 중 어떤 검사인지 자동으로 인식합니다."),
        ("자동 채점", "BBS 채점 기준에 따라 0-4점 자동 채점. 판단 근거도 함께 제공됩니다."),
        ("영상 동기화", "정면/측면 두 영상을 오디오 분석으로 자동 동기화합니다."),
        ("실시간 포즈 추적", "YOLOv8로 17개 관절을 실시간 추적하여 정밀한 분석을 수행합니다."),
        ("상세 분석 리포트", "좌우/앞뒤 흔들림, 관절 각도, 유지 시간 등 정량적 데이터 제공"),
        ("다중 인물 처리", "영상에 여러 사람이 있어도 환자를 자동으로 식별합니다.")
    ]
)

# Slide 12: AI Technology
add_content_slide(
    "AI 기술 스택",
    [
        "YOLOv8-pose: 실시간 멀티 포즈 추정 (17개 관절)",
        "Audio Cross-Correlation: FFT 기반 영상 동기화",
        "Rule-based Scoring: BBS 기준 기반 채점 엔진",
        "Stability Analysis: CoM 궤적, Sway 메트릭 분석",
        "Action Recognition: 동작 패턴 기반 검사 자동 인식"
    ]
)

# Slide 13: Measurement Details
add_feature_slide(
    "AI 측정 항목",
    [
        ("관절 각도", "무릎, 엉덩이, 발목, 어깨 등 주요 관절의 굴곡/신전 각도 측정"),
        ("무게중심 (CoM)", "신체 무게중심의 궤적을 분석하여 균형 안정성 평가"),
        ("좌우 흔들림", "Medial-Lateral 방향의 신체 동요를 cm 단위로 측정"),
        ("앞뒤 흔들림", "Anterior-Posterior 방향의 신체 동요를 cm 단위로 측정"),
        ("유지 시간", "특정 자세를 유지한 시간을 초 단위로 정밀 측정"),
        ("스텝 카운트", "발 교차 횟수, 회전 시 스텝 수 등 자동 카운팅")
    ]
)

# Slide 14: Section - 사용자 인터페이스
add_section_slide("사용자 인터페이스")

# Slide 15: UI Overview
add_feature_slide(
    "UI 구성",
    [
        ("메인 화면", "자동 분석 / 수동 평가 선택\n최근 평가 이력 표시"),
        ("자동 분석", "영상 업로드 → AI 분석 → 결과 확인\n원클릭 완전 자동화"),
        ("수동 평가", "14개 검사 중 선택\n개별 검사 분석"),
        ("결과 화면", "점수, 신뢰도, 판단 근거\n상세 측정값 표시"),
        ("듀얼 플레이어", "정면/측면 동기화 재생\n분석 구간 확인"),
        ("리포트", "PDF 내보내기\n환자별 이력 관리")
    ]
)

# Slide 16: Section - 기술 아키텍처
add_section_slide("기술 아키텍처")

# Slide 17: Tech Stack
add_comparison_slide(
    "기술 스택",
    "Frontend",
    [
        "React 18 + TypeScript",
        "Vite (빌드 도구)",
        "Tailwind CSS (스타일링)",
        "Axios (API 통신)",
        "React Router (라우팅)"
    ],
    "Backend",
    [
        "Python 3.11+ / FastAPI",
        "YOLOv8-pose (AI 포즈 추정)",
        "OpenCV (영상 처리)",
        "NumPy/SciPy (수치 연산)",
        "FFmpeg (오디오 추출)"
    ]
)

# Slide 18: System Requirements
add_table_slide(
    "시스템 요구사항",
    ["항목", "최소 사양", "권장 사양"],
    [
        ["CPU", "Intel i5 / AMD Ryzen 5", "Intel i7 / AMD Ryzen 7"],
        ["RAM", "8GB", "16GB"],
        ["GPU", "-", "NVIDIA RTX 3060+"],
        ["저장공간", "10GB", "50GB"],
        ["OS", "Windows 10, macOS 11, Ubuntu 20.04", "최신 버전"]
    ]
)

# Slide 19: Section - 타겟 시장
add_section_slide("타겟 시장")

# Slide 20: Target Users
add_feature_slide(
    "타겟 사용자",
    [
        ("물리치료사", "균형 평가 자동화로 업무 효율성 향상\n객관적 데이터 기반 치료 계획 수립"),
        ("재활의학과", "표준화된 평가로 의료진 간 소통 개선\n치료 효과 정량적 추적"),
        ("요양/재활병원", "다수 환자 평가 효율화\n낙상 위험 환자 조기 식별"),
        ("노인복지시설", "정기적 균형 능력 모니터링\n낙상 예방 프로그램 지원"),
        ("스포츠 재활센터", "선수 균형 능력 평가\n부상 후 복귀 시점 판단"),
        ("연구기관", "대규모 데이터 수집\n균형 능력 연구 지원")
    ]
)

# Slide 21: Section - 향후 계획
add_section_slide("향후 개발 계획")

# Slide 22: Roadmap
add_table_slide(
    "개발 로드맵",
    ["Phase", "기능", "예상 시기"],
    [
        ["Phase 2", "환자 정보 관리, 평가 이력 저장", "Q2 2024"],
        ["Phase 2", "PDF 리포트 생성, 다국어 지원", "Q2 2024"],
        ["Phase 3", "클라우드 배포 (AWS/GCP)", "Q3 2024"],
        ["Phase 3", "모바일 앱 (React Native)", "Q3 2024"],
        ["Phase 3", "EMR/HIS 시스템 연동", "Q4 2024"],
        ["AI 고도화", "딥러닝 기반 채점 모델 학습", "2025"]
    ]
)

# Slide 23: Section - 한계점 및 고려사항
add_section_slide("한계점 및 고려사항")

# Slide 24: Limitations - Technical
add_feature_slide(
    "기술적 한계점",
    [
        ("실시간 처리 불가", "현재는 녹화된 영상 분석만 지원하며, 실시간 카메라로 검사를 진행하는 것은 불가능합니다."),
        ("영상 품질 의존성", "저화질, 어두운 조명, 흔들리는 영상에서는 포즈 인식 정확도가 낮아질 수 있습니다."),
        ("카메라 각도", "정면/측면 촬영 각도가 정확하지 않으면 관절 각도 측정에 오차가 발생합니다."),
        ("의복 영향", "헐렁한 옷이나 긴 치마는 관절 포인트 감지를 방해할 수 있습니다."),
        ("다중 인물", "여러 사람이 동시에 있으면 환자 식별이 어려울 수 있습니다. (가장 큰 대상 선택)"),
        ("서버 필요", "현재 버전은 백엔드 서버가 필요하며, 오프라인 사용이 제한됩니다.")
    ]
)

# Slide 25: Limitations - Clinical
add_feature_slide(
    "임상적 고려사항",
    [
        ("보조 도구", "AI는 임상 전문가의 판단을 '대체'하는 것이 아닌 '보조'하는 도구입니다."),
        ("채점 정확도", "예상 정확도 85%+로, 일부 경계 사례에서 전문가 확인이 필요합니다."),
        ("표준 환경", "검사 환경(의자 높이, 발판 등)이 표준과 다르면 채점에 영향을 줄 수 있습니다."),
        ("실시간 중재 불가", "영상 기반 분석으로 검사 중 즉각적인 중재/보호가 불가능합니다."),
        ("학습 데이터 한계", "현재 Rule-based 방식으로, 다양한 임상 사례 학습이 제한적입니다."),
        ("의료기기 미인증", "현재 연구/보조 목적이며, 의료기기 인증 절차가 필요합니다.")
    ]
)

# Slide 26: Limitations with Solutions
add_comparison_slide(
    "한계점 및 해결 방안",
    "현재 한계",
    [
        "실시간 카메라 검사 불가 (녹화 영상만)",
        "영상 품질에 따른 정확도 변동",
        "Rule-based 채점의 유연성 부족",
        "오프라인 사용 불가",
        "모바일 앱 및 EMR 연동 미지원"
    ],
    "개선 계획",
    [
        "실시간 스트리밍 분석 기능 개발",
        "영상 전처리 및 품질 자동 보정",
        "딥러닝 기반 채점 모델 개발",
        "Edge AI / 로컬 처리 지원",
        "모바일 앱 및 HL7 FHIR EMR 연동"
    ]
)

# Slide 27: Best Practices
add_content_slide(
    "최적의 사용을 위한 권장사항",
    [
        "밝은 조명 환경에서 촬영 (자연광 또는 균일한 인공조명)",
        "카메라를 환자 정면/측면 90도 각도로 고정 설치",
        "환자는 몸에 맞는 옷 착용 권장 (반바지, 티셔츠)",
        "배경은 단색으로 하여 환자와 구분되도록 설정",
        "치료사는 가능한 프레임 밖에 위치",
        "검사 전 샘플 영상으로 포즈 인식 상태 확인",
        "AI 채점 결과는 전문가 검토 후 최종 확정"
    ]
)

# Slide 28: Closing
add_title_slide(
    "감사합니다",
    "AI로 더 정확하고 효율적인 균형 평가를"
)

# Save presentation
output_path = "/Users/aisoft/Documents/Camera/BBS_AI_Assessment_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
